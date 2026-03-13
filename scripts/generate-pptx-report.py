#!/usr/bin/env python3
"""
広告レポート PPTX — World-Class Edition
16:9 / Dark Modern / Design-First
"""

import json, os, subprocess, sys, tempfile, datetime, math

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import matplotlib.patheffects as pe
from matplotlib.colors import LinearSegmentedColormap
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

plt.rcParams["font.family"] = "Hiragino Sans"
plt.rcParams["axes.unicode_minus"] = False

# ═══ Design System ═══
W = Inches(13.333)
H = Inches(7.5)

class C:
    BG          = RGBColor(0x09, 0x09, 0x17)
    BG2         = RGBColor(0x0D, 0x11, 0x23)
    SURFACE     = RGBColor(0x12, 0x17, 0x2B)
    CARD        = RGBColor(0x16, 0x1D, 0x35)
    CARD_HOVER  = RGBColor(0x1C, 0x24, 0x40)
    BORDER      = RGBColor(0x1F, 0x2A, 0x48)
    BORDER_L    = RGBColor(0x2A, 0x3A, 0x5C)
    WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
    TEXT        = RGBColor(0xE2, 0xE8, 0xF0)
    TEXT2       = RGBColor(0x94, 0xA3, 0xB8)
    MUTED       = RGBColor(0x64, 0x74, 0x8B)
    # Accent
    BLUE        = RGBColor(0x38, 0xBD, 0xF8)
    INDIGO      = RGBColor(0x81, 0x8C, 0xF8)
    VIOLET      = RGBColor(0xA7, 0x8B, 0xFA)
    CYAN        = RGBColor(0x22, 0xD3, 0xEE)
    GREEN       = RGBColor(0x34, 0xD3, 0x99)
    EMERALD     = RGBColor(0x10, 0xB9, 0x81)
    AMBER       = RGBColor(0xFB, 0xBF, 0x24)
    ORANGE      = RGBColor(0xFB, 0x92, 0x3C)
    RED         = RGBColor(0xF8, 0x71, 0x71)
    ROSE        = RGBColor(0xFB, 0x7B, 0xA2)
    TEAL        = RGBColor(0x2D, 0xD4, 0xBF)
    PINK        = RGBColor(0xF4, 0x72, 0xB6)

M = {
    "bg": "#090917", "surface": "#12172B", "card": "#161D35",
    "border": "#1F2A48", "white": "#E2E8F0", "muted": "#64748B",
    "blue": "#38BDF8", "indigo": "#818CF8", "violet": "#A78BFA",
    "cyan": "#22D3EE", "green": "#34D399", "emerald": "#10B981",
    "amber": "#FBBF24", "orange": "#FB923C", "red": "#F87171",
    "rose": "#FB7BA2", "teal": "#2DD4BF", "pink": "#F472B6",
    "dark": "#0D1123", "text2": "#94A3B8",
}

# ═══ Helpers ═══

def fmt(v):
    if abs(v) >= 1_000_000: return f"¥{v/1_000_000:.1f}M"
    if abs(v) >= 1_000: return f"¥{v/1000:.0f}K"
    return f"¥{int(v):,}"

def fmt_full(v): return f"¥{int(v):,}"
def fmt_pct(v): return f"{v:.2f}%"
def fmt_num(v): return f"{int(v):,}"

def apply_fee(spend, rate, method):
    if method == "margin":
        return spend / (1 - rate / 100) if rate < 100 else spend
    return spend * (1 + rate / 100)


# ═══ XML helpers ═══

def _gradient(shape, c1, c2, angle=90):
    spPr = shape._element.spPr
    fill = parse_xml(f'''
    <a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
        <a:gsLst>
            <a:gs pos="0"><a:srgbClr val="{c1}"/></a:gs>
            <a:gs pos="100000"><a:srgbClr val="{c2}"/></a:gs>
        </a:gsLst>
        <a:lin ang="{angle * 60000}" scaled="1"/>
    </a:gradFill>''')
    for tag in ['a:solidFill', 'a:gradFill', 'a:noFill']:
        e = spPr.find(qn(tag))
        if e is not None: spPr.remove(e)
    spPr.insert(0, fill)

def _gradient3(shape, c1, c2, c3, angle=90):
    spPr = shape._element.spPr
    fill = parse_xml(f'''
    <a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
        <a:gsLst>
            <a:gs pos="0"><a:srgbClr val="{c1}"/></a:gs>
            <a:gs pos="50000"><a:srgbClr val="{c2}"/></a:gs>
            <a:gs pos="100000"><a:srgbClr val="{c3}"/></a:gs>
        </a:gsLst>
        <a:lin ang="{angle * 60000}" scaled="1"/>
    </a:gradFill>''')
    for tag in ['a:solidFill', 'a:gradFill', 'a:noFill']:
        e = spPr.find(qn(tag))
        if e is not None: spPr.remove(e)
    spPr.insert(0, fill)

def _shadow(shape, blur=50000, dist=18000, alpha=25):
    spPr = shape._element.spPr
    fx = parse_xml(f'''
    <a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <a:outerShdw blurRad="{blur}" dist="{dist}" dir="5400000" rotWithShape="0">
            <a:srgbClr val="000000"><a:alpha val="{alpha*1000}"/></a:srgbClr>
        </a:outerShdw>
    </a:effectLst>''')
    e = spPr.find(qn('a:effectLst'))
    if e is not None: spPr.remove(e)
    spPr.append(fx)

def _round(shape, r=60000):
    pg = shape._element.spPr.find(qn('a:prstGeom'))
    if pg is not None:
        av = pg.find(qn('a:avLst'))
        if av is None:
            av = parse_xml('<a:avLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
            pg.append(av)
        av.clear()
        av.append(parse_xml(f'<a:gd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="adj" fmla="val {r}"/>'))

def _glow(shape, color_hex, radius=120000, alpha=30):
    spPr = shape._element.spPr
    fx = spPr.find(qn('a:effectLst'))
    if fx is None:
        fx = parse_xml('<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
        spPr.append(fx)
    glow = parse_xml(f'''
    <a:glow xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rad="{radius}">
        <a:srgbClr val="{color_hex}"><a:alpha val="{alpha*1000}"/></a:srgbClr>
    </a:glow>''')
    fx.append(glow)

def _alpha(shape, alpha_pct):
    """図形の透明度を設定"""
    fill_elem = shape._element.spPr.find(qn('a:solidFill'))
    if fill_elem is not None:
        clr = fill_elem.find(qn('a:srgbClr'))
        if clr is not None:
            a = parse_xml(f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="{alpha_pct * 1000}"/>')
            clr.append(a)


# ═══ Layout Components ═══

def dark_bg(slide):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    _gradient(r, "090917", "0D1123", 135)
    r.line.fill.background()

def add_orb(slide, x, y, size, c1, c2, alpha=15):
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    _gradient(o, c1, c2, 135)
    o.line.fill.background()
    # alpha via XML
    spPr = o._element.spPr
    gf = spPr.find(qn('a:gradFill'))
    if gf is not None:
        for gs in gf.findall(qn('a:gs')):
            clr = gs.find(qn('a:srgbClr'))
            if clr is not None:
                a = parse_xml(f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="{alpha*1000}"/>')
                clr.append(a)

def add_text(slide, text, x, y, w, h, size=10, color=None, bold=False, align=PP_ALIGN.LEFT, spacing=1.15, font_name=None):
    if color is None: color = C.TEXT
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        if font_name:
            p.font.name = font_name
        # spacing
        pPr = p._pPr
        if pPr is None:
            pPr = parse_xml('<a:pPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
            p._p.insert(0, pPr)
        ls = parse_xml(f'<a:lnSpc xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:spcPct val="{int(spacing*100)}000"/></a:lnSpc>')
        old = pPr.find(qn('a:lnSpc'))
        if old is not None: pPr.remove(old)
        pPr.append(ls)
    return tb

def add_glass_card(slide, x, y, w, h, accent_hex=None):
    """Glassmorphism風カード"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = C.CARD
    card.line.color.rgb = C.BORDER
    card.line.width = Pt(0.5)
    _round(card, 40000)
    _shadow(card, blur=60000, dist=12000, alpha=20)
    # subtle alpha
    spPr = card._element.spPr
    sf = spPr.find(qn('a:solidFill'))
    if sf is not None:
        clr = sf.find(qn('a:srgbClr'))
        if clr is not None:
            a = parse_xml(f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="85000"/>')
            clr.append(a)
    if accent_hex:
        acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Emu(15000), y, Inches(0.04), h)
        acc.fill.solid()
        acc.fill.fore_color.rgb = RGBColor.from_string(accent_hex) if isinstance(accent_hex, str) else accent_hex
        acc.line.fill.background()
    return card

def page_num(slide, n, total):
    add_text(slide, f"{n:02d} — {total:02d}", Inches(12.2), Inches(7.1), Inches(1), Inches(0.3),
             size=7, color=C.MUTED, align=PP_ALIGN.RIGHT)

def section_line(slide, y=Inches(7.05)):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(12.1), Emu(8000))
    _gradient(ln, "1F2A48", "38BDF8", 0)
    ln.line.fill.background()


# ═══ Premium Table ═══

def add_table(slide, headers, rows, x, y, w, h, col_w=None):
    nr = len(rows) + 1
    nc = len(headers)
    ts = slide.shapes.add_table(nr, nc, x, y, w, h)
    tbl = ts.table
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = int(w * cw / total)
    for j, hdr in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = hdr
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0x16, 0x1D, 0x35)
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(7.5)
            p.font.bold = True
            p.font.color.rgb = C.BLUE
            p.alignment = PP_ALIGN.CENTER
        c.text_frame.margin_top = Pt(5)
        c.text_frame.margin_bottom = Pt(5)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i+1, j)
            c.text = str(val)
            c.fill.solid()
            c.fill.fore_color.rgb = C.SURFACE if i % 2 == 0 else C.BG2
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(7)
                p.font.color.rgb = C.TEXT
                p.alignment = PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT
            c.text_frame.margin_top = Pt(3)
            c.text_frame.margin_bottom = Pt(3)
            c.text_frame.margin_left = Pt(5)
            c.text_frame.margin_right = Pt(5)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
    # border styling via XML
    tblPr = tbl._tbl.tblPr
    if tblPr is None:
        tblPr = parse_xml('<a:tblPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
        tbl._tbl.insert(0, tblPr)
    tblPr.set('bandRow', '0')
    return ts


# ═══ matplotlib Dark Charts ═══

def _dark_style(ax, title=None):
    ax.set_facecolor(M["surface"])
    for sp in ax.spines.values():
        sp.set_visible(False)
    ax.tick_params(axis="both", colors=M["muted"], labelsize=7)
    ax.grid(axis="y", color=M["border"], linewidth=0.3, alpha=0.5, zorder=0)
    if title:
        ax.set_title(title, fontsize=11, fontweight="bold", color=M["white"], pad=10)

def chart_daily(data, path):
    daily = data.get("daily", [])
    if not daily: return None
    fr, fm = data.get("feeRate", 1), data.get("feeCalcMethod", "markup")
    dates = [d.get("date_start","")[-5:] for d in daily]
    spends = [apply_fee(d.get("spend",0), fr, fm)/1000 for d in daily]
    cvs = [d.get("cv",0) for d in daily]
    cpas = [d.get("cpa",0)/1000 for d in daily]

    fig, ax = plt.subplots(figsize=(11, 3.5), facecolor=M["dark"])
    x = np.arange(len(dates))
    # Area fill under bars
    ax.fill_between(x, 0, spends, alpha=0.15, color=M["blue"], zorder=1)
    bars = ax.bar(x, spends, color=M["blue"], alpha=0.7, width=0.5, zorder=2,
                  edgecolor=M["blue"], linewidth=0.5)
    # Gradient effect on bars
    for bar in bars:
        bar.set_alpha(0.75)

    _dark_style(ax)
    ax.set_ylabel("消化額 (K)", color=M["blue"], fontsize=8, labelpad=8)
    labels = dates if len(dates) <= 16 else [d if i % 2 == 0 else "" for i, d in enumerate(dates)]
    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=6.5, color=M["muted"], rotation=30 if len(dates) > 10 else 0)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: f"¥{v:.0f}K"))

    ax2 = ax.twinx()
    ax2.plot(x, cvs, color=M["green"], linewidth=2.5, marker="o", markersize=5,
             markerfacecolor=M["dark"], markeredgecolor=M["green"], markeredgewidth=2,
             zorder=4, path_effects=[pe.withStroke(linewidth=4, foreground=M["dark"])])
    ax2.set_ylabel("CV", color=M["green"], fontsize=8, labelpad=8)
    ax2.tick_params(axis="y", colors=M["green"], labelsize=7)
    for sp in ax2.spines.values(): sp.set_visible(False)

    for i, cv in enumerate(cvs):
        if len(dates) <= 16 or i % 2 == 0:
            ax2.annotate(str(int(cv)), (x[i], cv), textcoords="offset points", xytext=(0, 10),
                         ha="center", fontsize=6.5, color=M["green"], fontweight="bold",
                         path_effects=[pe.withStroke(linewidth=2, foreground=M["dark"])])

    # CPA line (thin)
    ax3 = ax.twinx()
    ax3.spines["right"].set_position(("outward", 40))
    ax3.plot(x, cpas, color=M["amber"], linewidth=1.5, linestyle="--", alpha=0.6, zorder=3)
    ax3.set_ylabel("CPA (K)", color=M["amber"], fontsize=7, labelpad=5)
    ax3.tick_params(axis="y", colors=M["amber"], labelsize=6)
    for sp in ax3.spines.values(): sp.set_visible(False)

    from matplotlib.lines import Line2D
    from matplotlib.patches import Patch
    legend = [
        Patch(facecolor=M["blue"], alpha=0.7, label="消化額"),
        Line2D([0],[0], color=M["green"], lw=2.5, marker="o", ms=5, label="CV"),
        Line2D([0],[0], color=M["amber"], lw=1.5, ls="--", label="CPA"),
    ]
    ax.legend(handles=legend, loc="upper left", fontsize=7, facecolor=M["surface"],
              edgecolor=M["border"], labelcolor=M["white"], framealpha=0.9)

    fig.tight_layout(pad=0.5)
    fig.savefig(path, dpi=300, bbox_inches="tight", facecolor=M["dark"])
    plt.close(fig)
    return path

def chart_campaigns(data, path):
    camps = data.get("campaigns", [])
    if not camps: return None
    top = sorted(camps, key=lambda c: c.get("spend",0), reverse=True)[:8]
    names = [c.get("campaign_name","")[:22] for c in top]
    spends = [c.get("spend",0)/1000 for c in top]
    cvs = [c.get("cv",0) for c in top]
    cpas = [c.get("cpa",0) for c in top]

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 3.0), facecolor=M["dark"],
                                   gridspec_kw={"width_ratios": [1.2, 1]})
    y = np.arange(len(names))

    # Spend bars
    colors_spend = [M["blue"]] * len(top)
    ax1.barh(y, spends, color=colors_spend, alpha=0.8, height=0.55, edgecolor=M["surface"], linewidth=0.5)
    ax1.set_yticks(y)
    ax1.set_yticklabels(names, fontsize=7, color=M["white"])
    ax1.invert_yaxis()
    _dark_style(ax1, "消化額 (¥K)")
    ax1.grid(axis="x", color=M["border"], linewidth=0.3, alpha=0.5)
    ax1.grid(axis="y", visible=False)
    for i, v in enumerate(spends):
        ax1.text(v + max(spends)*0.02, i, f"¥{v:.0f}K", va="center", fontsize=6.5, color=M["blue"], fontweight="bold")

    # CPA comparison
    cpa_colors = [M["green"] if c < 7000 else M["amber"] if c < 10000 else M["red"] for c in cpas]
    ax2.barh(y, cpas, color=cpa_colors, alpha=0.85, height=0.55)
    ax2.set_yticks(y)
    ax2.set_yticklabels([""] * len(y))
    ax2.invert_yaxis()
    _dark_style(ax2, "CPA (効率)")
    ax2.grid(axis="x", color=M["border"], linewidth=0.3, alpha=0.5)
    ax2.grid(axis="y", visible=False)
    for i, (v, cv) in enumerate(zip(cpas, cvs)):
        ax2.text(v + max(cpas)*0.02, i, f"¥{v:,.0f}  ({cv}CV)", va="center", fontsize=6.5,
                 color=cpa_colors[i], fontweight="bold")

    fig.tight_layout(pad=1.0)
    fig.savefig(path, dpi=300, bbox_inches="tight", facecolor=M["dark"])
    plt.close(fig)
    return path

def chart_device(data, path):
    devs = data.get("deviceBreakdown", [])
    if not devs: return None
    labels = [d.get("device","") for d in devs]
    spends = [d.get("spend",0) for d in devs]
    total = sum(spends)
    shares = [s/total*100 if total else 0 for s in spends]
    cvs = [d.get("cv",0) for d in devs]
    cpas = [d.get("cpa",0) for d in devs]

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 3.2), facecolor=M["dark"],
                                   gridspec_kw={"width_ratios": [1, 1.3]})
    colors = [M["blue"], M["violet"], M["teal"], M["pink"]][:len(devs)]

    # Donut
    wedges, _, autotexts = ax1.pie(shares, labels=None, autopct=lambda p: f"{p:.1f}%",
                                    startangle=90, colors=colors,
                                    pctdistance=0.75, wedgeprops={"width": 0.4, "edgecolor": M["dark"], "linewidth": 3})
    for t in autotexts:
        t.set_fontweight("bold")
        t.set_fontsize(10)
        t.set_color(M["white"])
    ax1.set_title("消化構成比", fontsize=11, fontweight="bold", color=M["white"], pad=12)
    ax1.legend(wedges, labels, loc="lower center", fontsize=8, ncol=3,
               facecolor=M["surface"], edgecolor=M["border"], labelcolor=M["white"],
               bbox_to_anchor=(0.5, -0.05))
    # Center text
    ax1.text(0, 0, f"¥{total/1000:.0f}K", ha="center", va="center",
             fontsize=14, fontweight="bold", color=M["white"])

    # Grouped bar
    x = np.arange(len(labels))
    w = 0.3
    b1 = ax2.bar(x - w/2, cvs, w, color=M["green"], alpha=0.85, label="CV", zorder=2)
    b2 = ax2.bar(x + w/2, [c/1000 for c in cpas], w, color=M["amber"], alpha=0.85, label="CPA (K)", zorder=2)
    ax2.set_xticks(x)
    ax2.set_xticklabels(labels, fontsize=9, color=M["white"])
    _dark_style(ax2, "CV × CPA")
    for bar, val in zip(b1, cvs):
        ax2.text(bar.get_x()+bar.get_width()/2, bar.get_height()+1, str(int(val)),
                 ha="center", fontsize=8, fontweight="bold", color=M["green"])
    for bar, val in zip(b2, cpas):
        ax2.text(bar.get_x()+bar.get_width()/2, val/1000+0.2, f"¥{val:,.0f}",
                 ha="center", fontsize=6.5, color=M["amber"])
    ax2.legend(fontsize=8, facecolor=M["surface"], edgecolor=M["border"], labelcolor=M["white"])

    fig.tight_layout(pad=0.6)
    fig.savefig(path, dpi=300, bbox_inches="tight", facecolor=M["dark"])
    plt.close(fig)
    return path

def chart_demo(data, path):
    demo = data.get("demographicBreakdown", [])
    if not demo: return None
    ag = {}
    for d in demo:
        a, g = d.get("age","?"), d.get("gender","?")
        if a not in ag: ag[a] = {}
        ag[a][g] = ag[a].get(g, 0) + d.get("cv", 0)
    ages = sorted(ag.keys())
    genders = sorted(set(g for a in ag.values() for g in a))

    fig, ax = plt.subplots(figsize=(7, 3.2), facecolor=M["dark"])
    x = np.arange(len(ages))
    w = 0.3
    gc = [M["blue"], M["rose"], M["teal"], M["amber"]]

    for gi, g in enumerate(genders):
        vals = [ag.get(a,{}).get(g,0) for a in ages]
        offset = (gi - (len(genders)-1)/2) * w
        bars = ax.bar(x+offset, vals, w*0.85, label=g, color=gc[gi%len(gc)], alpha=0.85, zorder=2,
                       edgecolor=M["dark"], linewidth=0.5)
        for bar, val in zip(bars, vals):
            if val > 0:
                ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.3,
                        str(int(val)), ha="center", fontsize=7, fontweight="bold", color=gc[gi%len(gc)])

    _dark_style(ax, "年齢 × 性別 CV分布")
    ax.set_xticks(x)
    ax.set_xticklabels(ages, fontsize=9, color=M["white"])
    ax.set_ylabel("CV", fontsize=8, color=M["white"])
    ax.legend(fontsize=8, facecolor=M["surface"], edgecolor=M["border"], labelcolor=M["white"])

    fig.tight_layout(pad=0.5)
    fig.savefig(path, dpi=300, bbox_inches="tight", facecolor=M["dark"])
    plt.close(fig)
    return path


# ═══ Slide Builder ═══

def build(data):
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    p = data.get("project", {})
    fr = data.get("feeRate", 1)
    fm = data.get("feeCalcMethod", "markup")
    sf = apply_fee(p.get("spend",0), fr, fm)
    cpc = p.get("spend",0) / p.get("clicks",1) if p.get("clicks",0) > 0 else 0
    name = data.get("projectName", "広告運用")
    daily = data.get("daily", [])
    dates = sorted(d.get("date_start","") for d in daily if d.get("date_start"))
    period = f"{dates[0]} — {dates[-1]}" if dates else ""
    today = datetime.date.today().strftime("%Y.%m.%d")
    analysis = data.get("analysis", {})
    overall = analysis.get("overall", {})
    client = analysis.get("clientReport", {})
    summary = client.get("summary") or overall.get("summary") or ""
    insights = overall.get("insights", [])
    improvements = client.get("improvements", overall.get("recommendations", []))
    camps = sorted(data.get("campaigns",[]), key=lambda x: x.get("spend",0), reverse=True)
    devs = data.get("deviceBreakdown", [])
    demo = data.get("demographicBreakdown", [])
    roas = p.get("purchase_roas")
    budget = data.get("monthlyBudget") or 0

    ts = 5
    if devs or demo: ts += 1
    if improvements: ts += 1
    sn = 0

    with tempfile.TemporaryDirectory() as tmp:

        # ══════ SLIDE 1: COVER ══════
        sn += 1
        s = prs.slides.add_slide(blank)
        dark_bg(s)
        # Orbs
        add_orb(s, Inches(8), Inches(-2), Inches(8), "1E3A8A", "3B82F6", 8)
        add_orb(s, Inches(-3), Inches(3), Inches(7), "581C87", "A78BFA", 6)
        add_orb(s, Inches(10), Inches(4), Inches(5), "064E3B", "34D399", 5)

        # Horizontal line
        ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.7), Inches(4), Emu(12000))
        _gradient(ln, "38BDF8", "818CF8", 0)
        ln.line.fill.background()

        # ADVERTISING REPORT
        add_text(s, "ADVERTISING", Inches(1), Inches(1.6), Inches(6), Inches(0.4),
                 size=13, color=C.BLUE, bold=False, font_name="Helvetica Neue")
        add_text(s, "PERFORMANCE\nREPORT", Inches(1), Inches(2.85), Inches(8), Inches(1.2),
                 size=44, color=C.WHITE, bold=True, spacing=1.0)
        add_text(s, name, Inches(1), Inches(4.2), Inches(8), Inches(0.5),
                 size=18, color=C.BLUE)
        add_text(s, period, Inches(1), Inches(4.75), Inches(4), Inches(0.3),
                 size=11, color=C.TEXT2)
        add_text(s, today, Inches(1), Inches(5.1), Inches(3), Inches(0.3),
                 size=10, color=C.MUTED)

        # Right side KPI preview cards
        kpi_preview = [
            ("TOTAL SPEND", fmt_full(sf), "38BDF8"),
            ("CONVERSIONS", f"{int(p.get('cv',0))}", "34D399"),
            ("CPA", fmt_full(p.get("cpa",0)), "FBBF24"),
            ("ROAS", f"{roas:.2f}x" if isinstance(roas,(int,float)) and roas > 0 else "—", "A78BFA"),
        ]
        for i, (lbl, val, clr) in enumerate(kpi_preview):
            cy = Inches(1.5 + i * 1.35)
            card = add_glass_card(s, Inches(9.5), cy, Inches(3.2), Inches(1.1), clr)
            add_text(s, lbl, Inches(9.8), cy + Inches(0.15), Inches(2.5), Inches(0.2),
                     size=8, color=C.MUTED)
            add_text(s, val, Inches(9.8), cy + Inches(0.4), Inches(2.8), Inches(0.5),
                     size=24, color=C.WHITE, bold=True)

        # Bottom line
        bln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.15), W, Emu(12000))
        _gradient3(bln, "38BDF8", "818CF8", "A78BFA", 0)
        bln.line.fill.background()
        add_text(s, "CONFIDENTIAL", Inches(0.5), Inches(7.2), Inches(3), Inches(0.25),
                 size=7, color=C.MUTED)

        # ══════ SLIDE 2: KPI DASHBOARD ══════
        sn += 1
        s = prs.slides.add_slide(blank)
        dark_bg(s)
        add_orb(s, Inches(9), Inches(-1), Inches(6), "1E3A8A", "3B82F6", 5)
        section_line(s)
        page_num(s, sn, ts)

        # Section title
        add_text(s, "PERFORMANCE", Inches(0.6), Inches(0.3), Inches(4), Inches(0.3),
                 size=10, color=C.BLUE, bold=True, font_name="Helvetica Neue")
        add_text(s, "パフォーマンスサマリー", Inches(0.6), Inches(0.6), Inches(6), Inches(0.4),
                 size=22, color=C.WHITE, bold=True)

        # KPI Cards row
        kpi_data = [
            ("消化額（Fee込）", fmt_full(sf), f"税前 {fmt_full(p.get('spend',0))}", "38BDF8", C.BLUE),
            ("CV数", f'{int(p.get("cv",0))}件', f"IMP {fmt_num(p.get('impressions',0))}", "34D399", C.GREEN),
            ("CPA", fmt_full(p.get("cpa",0)) if p.get("cv",0)>0 else "—", f"目標 ¥10,000", "FBBF24", C.AMBER),
            ("CTR", fmt_pct(p.get("ctr",0)), f"Click {fmt_num(p.get('clicks',0))}", "818CF8", C.INDIGO),
            ("CPC", fmt_full(cpc), "", "22D3EE", C.CYAN),
        ]
        if isinstance(roas,(int,float)) and roas > 0:
            kpi_data.append(("ROAS", f"{roas:.2f}x", "", "A78BFA", C.VIOLET))

        n = len(kpi_data)
        cw = Inches(1.85)
        gap = Inches(0.15)
        total_w = n * cw + (n-1) * gap
        sx = int((W - total_w) / 2)

        for i, (lbl, val, sub, clr_hex, clr) in enumerate(kpi_data):
            x = sx + i * int(cw + gap)
            card = add_glass_card(s, x, Inches(1.2), cw, Inches(1.25), clr_hex)
            add_text(s, lbl, x + Inches(0.15), Inches(1.35), cw - Inches(0.3), Inches(0.2),
                     size=8, color=C.MUTED, align=PP_ALIGN.CENTER)
            add_text(s, val, x + Inches(0.05), Inches(1.6), cw - Inches(0.1), Inches(0.45),
                     size=22, color=C.WHITE, bold=True, align=PP_ALIGN.CENTER)
            if sub:
                add_text(s, sub, x + Inches(0.1), Inches(2.1), cw - Inches(0.2), Inches(0.2),
                         size=7, color=C.TEXT2, align=PP_ALIGN.CENTER)

        # Budget bar
        if budget > 0:
            rate = sf / budget * 100
            add_text(s, f"月間予算  {fmt_full(budget)}", Inches(0.8), Inches(2.7), Inches(3), Inches(0.25),
                     size=9, color=C.TEXT2)
            # Bar bg
            bbg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.0), Inches(11.7), Inches(0.2))
            bbg.fill.solid()
            bbg.fill.fore_color.rgb = C.SURFACE
            bbg.line.fill.background()
            _round(bbg, 50000)
            # Bar fg
            bw = min(Inches(11.7) * rate / 100, Inches(11.7))
            if bw > 0:
                bfg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.0), int(bw), Inches(0.2))
                c1 = "38BDF8" if rate <= 50 else "FBBF24" if rate <= 80 else "F87171"
                c2 = "818CF8" if rate <= 50 else "FB923C" if rate <= 80 else "EF4444"
                _gradient(bfg, c1, c2, 0)
                bfg.line.fill.background()
                _round(bfg, 50000)

            dom = datetime.date.today().day
            est = int(sf / max(dom / 31, 0.01))
            add_text(s, f"消化 {fmt_full(sf)} ({rate:.1f}%)    |    着地予想 {fmt_full(est)}",
                     Inches(0.8), Inches(3.25), Inches(8), Inches(0.2), size=8, color=C.TEXT2)

        # Summary cards
        if summary:
            add_glass_card(s, Inches(0.6), Inches(3.7), Inches(6), Inches(1.2), "38BDF8")
            add_text(s, "📊 REPORT", Inches(0.85), Inches(3.8), Inches(2), Inches(0.2),
                     size=8, color=C.BLUE, bold=True)
            rpt = f"消化額 {fmt_full(sf)} に対し CV {int(p.get('cv',0))}件。CPA {fmt_full(p.get('cpa',0))} で目標内推移。"
            if isinstance(roas,(int,float)) and roas > 0:
                rpt += f" ROAS {roas:.2f}x。"
            add_text(s, rpt, Inches(0.85), Inches(4.05), Inches(5.5), Inches(0.7),
                     size=9, color=C.TEXT, spacing=1.4)

            add_glass_card(s, Inches(6.8), Inches(3.7), Inches(6), Inches(1.2), "34D399")
            add_text(s, "💡 INSIGHT", Inches(7.05), Inches(3.8), Inches(2), Inches(0.2),
                     size=8, color=C.GREEN, bold=True)
            add_text(s, summary, Inches(7.05), Inches(4.05), Inches(5.5), Inches(0.7),
                     size=9, color=C.TEXT, spacing=1.4)

        # Weekly mini table
        if daily:
            weeks = {}
            for d in daily:
                try:
                    dt = datetime.date.fromisoformat(d.get("date_start",""))
                    wk = f"W{dt.isocalendar()[1]}"
                except: wk = "?"
                if wk not in weeks: weeks[wk] = {"spend":0,"cv":0,"imp":0,"click":0}
                weeks[wk]["spend"] += d.get("spend",0)
                weeks[wk]["cv"] += d.get("cv",0)
                weeks[wk]["imp"] += d.get("impressions",0)
                weeks[wk]["click"] += d.get("clicks",0)
            hdrs = ["週", "消化額", "CV", "CPA", "CTR"]
            rows = []
            for wk, wd in sorted(weeks.items()):
                sp = apply_fee(wd["spend"], fr, fm)
                cpa = sp / wd["cv"] if wd["cv"] > 0 else 0
                ctr = wd["click"] / wd["imp"] * 100 if wd["imp"] > 0 else 0
                rows.append([wk, fmt_full(sp), str(wd["cv"]), fmt_full(cpa), fmt_pct(ctr)])
            if rows:
                add_table(s, hdrs, rows, Inches(0.6), Inches(5.2), Inches(12.1), Inches(0.2+0.2*len(rows)),
                          col_w=[0.5,1.2,0.4,1,0.6])

        # ══════ SLIDE 3: DAILY ══════
        sn += 1
        s = prs.slides.add_slide(blank)
        dark_bg(s)
        section_line(s)
        page_num(s, sn, ts)
        add_text(s, "DAILY TREND", Inches(0.6), Inches(0.25), Inches(4), Inches(0.25),
                 size=10, color=C.BLUE, bold=True, font_name="Helvetica Neue")
        add_text(s, "日次パフォーマンス推移", Inches(0.6), Inches(0.5), Inches(6), Inches(0.35),
                 size=20, color=C.WHITE, bold=True)

        cp = os.path.join(tmp, "daily.png")
        if chart_daily(data, cp):
            s.shapes.add_picture(cp, Inches(0.3), Inches(1.0), Inches(12.7), Inches(3.3))

        recent = sorted(daily, key=lambda x: x.get("date_start",""))[-10:]
        if recent:
            hdrs = ["日付", "消化額", "Fee込", "IMP", "Click", "CTR", "CV", "CPA"]
            rows = []
            for d in recent:
                sp = d.get("spend",0)
                rows.append([
                    d.get("date_start","")[-5:],
                    fmt_full(sp), fmt_full(apply_fee(sp, fr, fm)),
                    fmt_num(d.get("impressions",0)), fmt_num(d.get("clicks",0)),
                    fmt_pct(d.get("ctr",0)), str(int(d.get("cv",0))),
                    fmt_full(d.get("cpa",0)) if d.get("cv",0) > 0 else "—"
                ])
            add_table(s, hdrs, rows, Inches(0.3), Inches(4.4), Inches(12.7),
                      Inches(0.2+0.18*len(rows)),
                      col_w=[0.7,1,1,0.9,0.7,0.6,0.4,0.9])

        # ══════ SLIDE 4: CAMPAIGNS ══════
        sn += 1
        s = prs.slides.add_slide(blank)
        dark_bg(s)
        section_line(s)
        page_num(s, sn, ts)
        add_text(s, "CAMPAIGNS", Inches(0.6), Inches(0.25), Inches(4), Inches(0.25),
                 size=10, color=C.BLUE, bold=True, font_name="Helvetica Neue")
        add_text(s, "キャンペーン別パフォーマンス", Inches(0.6), Inches(0.5), Inches(8), Inches(0.35),
                 size=20, color=C.WHITE, bold=True)

        top10 = camps[:10]
        if top10:
            total_sp = sum(c.get("spend",0) for c in top10)
            hdrs = ["キャンペーン", "消化額", "IMP", "Click", "CTR", "CV", "CPA", "構成比"]
            rows = []
            for c in top10:
                sh = c.get("spend",0)/total_sp*100 if total_sp else 0
                rows.append([
                    c.get("campaign_name","")[:30], fmt_full(c.get("spend",0)),
                    fmt_num(c.get("impressions",0)), fmt_num(c.get("clicks",0)),
                    fmt_pct(c.get("ctr",0)), str(int(c.get("cv",0))),
                    fmt_full(c.get("cpa",0)) if c.get("cv",0)>0 else "—", fmt_pct(sh)
                ])
            add_table(s, hdrs, rows, Inches(0.3), Inches(1.0), Inches(12.7),
                      Inches(0.2+0.19*len(rows)), col_w=[3,1,0.9,0.7,0.6,0.4,0.9,0.6])

        cp = os.path.join(tmp, "camps.png")
        if chart_campaigns(data, cp):
            cy = Inches(1.0 + 0.19*(len(top10)+1) + 0.1)
            cy = min(cy, Inches(4.0))
            remain = Inches(6.9) - cy
            if remain > Inches(1.5):
                s.shapes.add_picture(cp, Inches(0.3), cy, Inches(12.7), min(remain, Inches(2.8)))

        # ══════ SLIDE 5: DEVICE & DEMO ══════
        if devs or demo:
            sn += 1
            s = prs.slides.add_slide(blank)
            dark_bg(s)
            section_line(s)
            page_num(s, sn, ts)
            add_text(s, "AUDIENCE", Inches(0.6), Inches(0.25), Inches(4), Inches(0.25),
                     size=10, color=C.BLUE, bold=True, font_name="Helvetica Neue")
            add_text(s, "デバイス・属性別分析", Inches(0.6), Inches(0.5), Inches(8), Inches(0.35),
                     size=20, color=C.WHITE, bold=True)

            if devs:
                dp = os.path.join(tmp, "device.png")
                if chart_device(data, dp):
                    s.shapes.add_picture(dp, Inches(0.3), Inches(1.0), Inches(12.7), Inches(3.0))

            if demo:
                dmp = os.path.join(tmp, "demo.png")
                if chart_demo(data, dmp):
                    s.shapes.add_picture(dmp, Inches(0.3), Inches(4.1), Inches(7), Inches(2.7))

                # Demo summary table
                age_sum = {}
                for d in demo:
                    a = d.get("age","?")
                    if a not in age_sum: age_sum[a] = {"spend":0,"cv":0}
                    age_sum[a]["spend"] += d.get("spend",0)
                    age_sum[a]["cv"] += d.get("cv",0)
                hdrs = ["年齢層", "消化額", "CV", "CPA"]
                rows = []
                for a in sorted(age_sum):
                    v = age_sum[a]
                    cpa = v["spend"]/v["cv"] if v["cv"]>0 else 0
                    rows.append([a, fmt_full(v["spend"]), str(v["cv"]), fmt_full(cpa)])
                add_table(s, hdrs, rows, Inches(7.5), Inches(4.1), Inches(5.5),
                          Inches(0.2+0.19*len(rows)), col_w=[0.8,1,0.5,1])

        # ══════ SLIDE 6: INSIGHTS ══════
        sn += 1
        s = prs.slides.add_slide(blank)
        dark_bg(s)
        add_orb(s, Inches(9), Inches(2), Inches(6), "581C87", "A78BFA", 4)
        section_line(s)
        page_num(s, sn, ts)
        add_text(s, "ANALYSIS", Inches(0.6), Inches(0.25), Inches(4), Inches(0.25),
                 size=10, color=C.BLUE, bold=True, font_name="Helvetica Neue")
        add_text(s, "分析・総評", Inches(0.6), Inches(0.5), Inches(8), Inches(0.35),
                 size=20, color=C.WHITE, bold=True)

        if summary:
            card = add_glass_card(s, Inches(0.6), Inches(1.1), Inches(12.1), Inches(1.0), "38BDF8")
            add_text(s, summary, Inches(0.95), Inches(1.2), Inches(11.5), Inches(0.8),
                     size=12, color=C.TEXT, spacing=1.5)

        if insights:
            add_text(s, "KEY INSIGHTS", Inches(0.6), Inches(2.3), Inches(3), Inches(0.25),
                     size=9, color=C.BLUE, bold=True)
            for i, ins in enumerate(insights[:6]):
                y = Inches(2.7 + i * 0.65)
                card = add_glass_card(s, Inches(0.6), y, Inches(12.1), Inches(0.55))
                # Number badge
                badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           Inches(0.8), y + Emu(45000), Inches(0.3), Inches(0.3))
                badge_colors = ["38BDF8", "34D399", "FBBF24", "A78BFA", "FB923C", "22D3EE"]
                badge.fill.solid()
                badge.fill.fore_color.rgb = RGBColor.from_string(badge_colors[i % len(badge_colors)])
                badge.line.fill.background()
                _round(badge, 30000)
                tf = badge.text_frame
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.margin_left = Pt(0)
                tf.margin_right = Pt(0)
                pp = tf.paragraphs[0]
                pp.text = str(i+1)
                pp.font.size = Pt(10)
                pp.font.bold = True
                pp.font.color.rgb = C.BG
                pp.alignment = PP_ALIGN.CENTER

                add_text(s, ins, Inches(1.25), y + Emu(35000), Inches(11.2), Inches(0.4),
                         size=10.5, color=C.TEXT, spacing=1.2)

        # ══════ SLIDE 7: ACTIONS ══════
        if improvements:
            sn += 1
            s = prs.slides.add_slide(blank)
            dark_bg(s)
            add_orb(s, Inches(-2), Inches(2), Inches(6), "064E3B", "34D399", 5)
            section_line(s)
            page_num(s, sn, ts)
            add_text(s, "NEXT ACTIONS", Inches(0.6), Inches(0.25), Inches(4), Inches(0.25),
                     size=10, color=C.GREEN, bold=True, font_name="Helvetica Neue")
            add_text(s, "改善施策・ネクストアクション", Inches(0.6), Inches(0.5), Inches(8), Inches(0.35),
                     size=20, color=C.WHITE, bold=True)

            accent_colors = ["F87171", "FBBF24", "38BDF8", "34D399", "A78BFA", "FB923C", "22D3EE", "F472B6"]
            for i, item in enumerate(improvements[:8]):
                y = Inches(1.15 + i * 0.72)
                card = add_glass_card(s, Inches(0.6), y, Inches(12.1), Inches(0.6), accent_colors[i % len(accent_colors)])
                # Badge
                badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           Inches(0.85), y + Emu(55000), Inches(0.32), Inches(0.32))
                badge.fill.solid()
                badge.fill.fore_color.rgb = RGBColor.from_string(accent_colors[i % len(accent_colors)])
                badge.line.fill.background()
                _round(badge, 25000)
                tf = badge.text_frame
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.margin_left = Pt(0); tf.margin_right = Pt(0)
                pp = tf.paragraphs[0]
                pp.text = str(i+1)
                pp.font.size = Pt(11)
                pp.font.bold = True
                pp.font.color.rgb = C.BG
                pp.alignment = PP_ALIGN.CENTER

                add_text(s, item, Inches(1.35), y + Emu(50000), Inches(11), Inches(0.42),
                         size=11, color=C.TEXT, spacing=1.25)

    return prs


def upload(path, title):
    cmd = ["gws", "drive", "files", "create",
           "--params", json.dumps({"uploadType":"multipart","fields":"id,webViewLink"}),
           "--json", json.dumps({"name":title,"mimeType":"application/vnd.google-apps.presentation"}),
           "--upload", path]
    r = subprocess.run(cmd, capture_output=True, text=True)
    if r.returncode != 0: return None, None
    d = json.loads(r.stdout)
    fid = d.get("id")
    link = d.get("webViewLink")
    if fid:
        subprocess.run(["gws","drive","permissions","create",
                        "--params",json.dumps({"fileId":fid}),
                        "--json",json.dumps({"role":"writer","type":"anyone"})],
                       capture_output=True, text=True)
    url = link or (f"https://docs.google.com/presentation/d/{fid}/edit" if fid else None)
    return fid, url


def main():
    if "--sample" in sys.argv:
        data = {
            "projectName": "株式会社サンプル",
            "datePreset": "last_30d",
            "project": {"spend":1250000,"impressions":850000,"clicks":12500,"ctr":1.47,"cv":185,"cpa":8108,"purchase_roas":3.24},
            "campaigns": [
                {"campaign_name":"Search_ブランドKW","spend":320000,"impressions":180000,"clicks":4500,"ctr":2.50,"cv":62,"cpa":5161},
                {"campaign_name":"Search_一般KW","spend":280000,"impressions":220000,"clicks":3300,"ctr":1.50,"cv":38,"cpa":7368},
                {"campaign_name":"Display_リターゲティング","spend":250000,"impressions":150000,"clicks":2100,"ctr":1.40,"cv":45,"cpa":5556},
                {"campaign_name":"Display_類似オーディエンス","spend":200000,"impressions":180000,"clicks":1500,"ctr":0.83,"cv":22,"cpa":9091},
                {"campaign_name":"Instagram_リール広告","spend":120000,"impressions":80000,"clicks":700,"ctr":0.88,"cv":12,"cpa":10000},
                {"campaign_name":"Facebook_フィード広告","spend":80000,"impressions":40000,"clicks":400,"ctr":1.00,"cv":6,"cpa":13333},
            ],
            "creatives": [],
            "daily": [
                {"date_start":"2026-03-01","spend":95000,"impressions":65000,"clicks":975,"ctr":1.5,"cv":14,"cpa":6786},
                {"date_start":"2026-03-02","spend":88000,"impressions":62000,"clicks":806,"ctr":1.3,"cv":12,"cpa":7333},
                {"date_start":"2026-03-03","spend":102000,"impressions":68000,"clicks":1088,"ctr":1.6,"cv":16,"cpa":6375},
                {"date_start":"2026-03-04","spend":110000,"impressions":72000,"clicks":1224,"ctr":1.7,"cv":18,"cpa":6111},
                {"date_start":"2026-03-05","spend":98000,"impressions":64000,"clicks":896,"ctr":1.4,"cv":13,"cpa":7538},
                {"date_start":"2026-03-06","spend":105000,"impressions":67000,"clicks":1005,"ctr":1.5,"cv":15,"cpa":7000},
                {"date_start":"2026-03-07","spend":92000,"impressions":60000,"clicks":780,"ctr":1.3,"cv":11,"cpa":8364},
                {"date_start":"2026-03-08","spend":115000,"impressions":74000,"clicks":1332,"ctr":1.8,"cv":19,"cpa":6053},
                {"date_start":"2026-03-09","spend":88000,"impressions":63000,"clicks":882,"ctr":1.4,"cv":13,"cpa":6769},
                {"date_start":"2026-03-10","spend":97000,"impressions":65000,"clicks":975,"ctr":1.5,"cv":14,"cpa":6929},
                {"date_start":"2026-03-11","spend":108000,"impressions":70000,"clicks":1120,"ctr":1.6,"cv":17,"cpa":6353},
                {"date_start":"2026-03-12","spend":78000,"impressions":58000,"clicks":812,"ctr":1.4,"cv":12,"cpa":6500},
                {"date_start":"2026-03-13","spend":74000,"impressions":62000,"clicks":806,"ctr":1.3,"cv":11,"cpa":6727},
            ],
            "analysis": {
                "overall": {
                    "summary":"3月前半は順調に推移。CPAは目標¥10,000を大幅に下回る¥8,108で効率良好。ROAS 3.24xと投資対効果も高水準を維持。",
                    "insights": [
                        "Search_ブランドKWが最もCPA効率が良く（¥5,161）、全体CVの33.5%を占める主力キャンペーン",
                        "Display_リターゲティングもCPA ¥5,556と高効率。リタゲリストの拡充で更なるCV増が見込める",
                        "SNS広告（Instagram/Facebook）はCPA高めで改善余地。クリエイティブA/Bテストで効率化を推進",
                        "週末（3/8, 3/9）のCV効率が高い傾向。週末の予算配分強化でパフォーマンス向上を狙う",
                    ],
                },
                "clientReport": {
                    "summary":"3月前半は月間予算¥3,000,000に対し消化率50.0%で計画通り。全体CPA ¥8,108は目標を大きく下回り好調。",
                    "improvements": [
                        "Facebook広告: CPA ¥13,333 → クリエイティブ刷新（動画 vs カルーセル A/Bテスト実施）",
                        "Display_類似: ターゲティング精度向上 → オーディエンスソース見直しとLookalike拡張テスト",
                        "Search_一般KW: 除外キーワード30件追加 → 無駄クリック15%削減目標",
                        "週末予算配分を平日比+15%に調整 → 効率の良い時間帯（10-14時, 20-23時）でCV最大化",
                    ],
                },
            },
            "feeRate":20,"feeCalcMethod":"markup","monthlyBudget":3000000,
            "deviceBreakdown": [
                {"device":"モバイル","spend":750000,"impressions":510000,"clicks":7500,"cv":115,"cpa":6522,"ctr":1.47},
                {"device":"デスクトップ","spend":375000,"impressions":255000,"clicks":3750,"cv":52,"cpa":7212,"ctr":1.47},
                {"device":"タブレット","spend":125000,"impressions":85000,"clicks":1250,"cv":18,"cpa":6944,"ctr":1.47},
            ],
            "demographicBreakdown": [
                {"age":"18-24","gender":"male","spend":40000,"impressions":30000,"clicks":500,"cv":8,"cpa":5000},
                {"age":"18-24","gender":"female","spend":50000,"impressions":35000,"clicks":600,"cv":12,"cpa":4167},
                {"age":"25-34","gender":"male","spend":100000,"impressions":70000,"clicks":1200,"cv":25,"cpa":4000},
                {"age":"25-34","gender":"female","spend":120000,"impressions":80000,"clicks":1400,"cv":35,"cpa":3429},
                {"age":"35-44","gender":"male","spend":110000,"impressions":75000,"clicks":1300,"cv":30,"cpa":3667},
                {"age":"35-44","gender":"female","spend":100000,"impressions":70000,"clicks":1200,"cv":28,"cpa":3571},
                {"age":"45-54","gender":"male","spend":80000,"impressions":55000,"clicks":900,"cv":18,"cpa":4444},
                {"age":"45-54","gender":"female","spend":70000,"impressions":50000,"clicks":800,"cv":15,"cpa":4667},
                {"age":"55-64","gender":"male","spend":40000,"impressions":30000,"clicks":400,"cv":8,"cpa":5000},
                {"age":"55-64","gender":"female","spend":30000,"impressions":25000,"clicks":300,"cv":6,"cpa":5000},
            ],
        }
    else:
        data = json.loads(sys.stdin.read())

    prs = build(data)
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        pptx_path = f.name
        prs.save(pptx_path)

    title = f'{data.get("projectName","広告運用")} Ad Report {datetime.date.today().isoformat()}.pptx'
    fid, url = upload(pptx_path, title)
    os.unlink(pptx_path)

    if url:
        print(json.dumps({"ok":True,"presentationId":fid,"presentationUrl":url}))
    else:
        print(json.dumps({"ok":False,"error":"Google Driveへのアップロードに失敗"}))
        sys.exit(1)

if __name__ == "__main__":
    main()
